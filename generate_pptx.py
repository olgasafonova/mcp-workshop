#!/usr/bin/env python3
"""
Generate editable PPTX from MCP Workshop HTML slides.
Uses python-pptx to create native PowerPoint with editable text/shapes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Brand colors from the HTML
COLORS = {
    "red": RGBColor(0xE5, 0x39, 0x35),       # #E53935
    "dark": RGBColor(0x1A, 0x1A, 0x1A),      # #1A1A1A
    "light": RGBColor(0xFA, 0xFA, 0xFA),     # #FAFAFA
    "gray": RGBColor(0x88, 0x88, 0x88),      # #888888
    "text": RGBColor(0x44, 0x44, 0x44),      # #444444
    "text_dark": RGBColor(0x33, 0x33, 0x33), # #333333
    "muted": RGBColor(0xAA, 0xAA, 0xAA),     # #AAAAAA
    "border": RGBColor(0xE0, 0xE0, 0xE0),    # #E0E0E0
    "code_bg": RGBColor(0xF0, 0xF0, 0xF0),   # #F0F0F0
    "green": RGBColor(0x7C, 0xB3, 0x42),     # #7CB342 (for code strings)
}

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def set_shape_fill(shape, color):
    """Set solid fill color for a shape."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def set_shape_border(shape, color, width=Pt(2)):
    """Set border/line for a shape."""
    shape.line.color.rgb = color
    shape.line.width = width

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=None, font_name="Arial", align=PP_ALIGN.LEFT):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    if color:
        p.font.color.rgb = color
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=20, color=None):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"→ {item}"
        p.font.size = Pt(font_size)
        p.font.name = "Arial"
        if color:
            p.font.color.rgb = color
        p.space_after = Pt(8)
    return txBox


def create_title_slide(prs):
    """Slide 01: Title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Light background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["light"])
    background.line.fill.background()

    # Tier badge
    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(3), Inches(0.4),
                 "WORKSHOP / 1 HOUR", font_size=12, bold=True, color=COLORS["gray"])

    # Main title
    add_text_box(slide, Inches(0.6), Inches(2), Inches(10), Inches(2),
                 "From Explorer\nto Architect", font_size=72, bold=True, color=COLORS["dark"])

    # Subtitle
    add_text_box(slide, Inches(0.6), Inches(4.5), Inches(9), Inches(1),
                 "Model Context Protocol: the standard that lets AI assistants take action in your apps",
                 font_size=24, color=COLORS["text"])

    # Footer
    add_text_box(slide, Inches(0.6), Inches(6.5), Inches(3), Inches(0.5),
                 "Olga Safonova", font_size=16, bold=True, color=COLORS["dark"])
    add_text_box(slide, Inches(8), Inches(6.5), Inches(4.5), Inches(0.7),
                 "Women in Product Nordics\nTechWomen Copenhagen", font_size=14,
                 color=COLORS["text"], align=PP_ALIGN.RIGHT)


def create_problem_slide(prs):
    """Slide 02: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["light"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(3), Inches(0.3),
                 "TIER 1 / EXPLORER", font_size=12, bold=True, color=COLORS["gray"])

    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(3), Inches(0.4),
                 "THE PROBLEM", font_size=12, bold=True, color=COLORS["red"])

    # Main quote
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI that knows everything and does nothing is an "
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.name = "Arial"

    run = p.add_run()
    run.text = "expensive search bar."
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = COLORS["red"]

    add_text_box(slide, Inches(0.6), Inches(4.5), Inches(9), Inches(1),
                 "Current AI can answer questions and generate text.\nBut it can't create a Miro board or update your roadmap.",
                 font_size=20, color=COLORS["text"])


def create_rag_vs_mcp_slide(prs):
    """Slide 03: RAG vs MCP"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["light"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(3), Inches(0.3),
                 "TIER 1 / EXPLORER", font_size=12, bold=True, color=COLORS["gray"])

    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(3), Inches(0.4),
                 "THE DISTINCTION", font_size=12, bold=True, color=COLORS["red"])

    add_text_box(slide, Inches(0.6), Inches(1.7), Inches(6), Inches(0.8),
                 "RAG vs MCP", font_size=48, bold=True, color=COLORS["dark"])

    # RAG column
    rag_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.8), Inches(0.15), Inches(3.5))
    set_shape_fill(rag_box, COLORS["gray"])
    rag_box.line.fill.background()

    add_text_box(slide, Inches(1), Inches(2.8), Inches(5), Inches(0.6),
                 "RAG", font_size=36, bold=True, color=COLORS["gray"])
    add_text_box(slide, Inches(1), Inches(3.4), Inches(5), Inches(0.5),
                 "Information flows IN", font_size=24, bold=True, color=COLORS["dark"])
    add_text_box(slide, Inches(1), Inches(4), Inches(5), Inches(2),
                 "Retrieval-Augmented Generation.\nAI pulls data from documents,\ndatabases, knowledge bases.\n\nResult: Better answers.",
                 font_size=18, color=COLORS["text"])

    # MCP column
    mcp_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.8), Inches(0.15), Inches(3.5))
    set_shape_fill(mcp_box, COLORS["red"])
    mcp_box.line.fill.background()

    add_text_box(slide, Inches(7.2), Inches(2.8), Inches(5), Inches(0.6),
                 "MCP", font_size=36, bold=True, color=COLORS["red"])
    add_text_box(slide, Inches(7.2), Inches(3.4), Inches(5), Inches(0.5),
                 "Commands flow OUT", font_size=24, bold=True, color=COLORS["dark"])
    add_text_box(slide, Inches(7.2), Inches(4), Inches(5), Inches(2),
                 "Model Context Protocol.\nAI sends actions to apps,\nAPIs, local software.\n\nResult: Things get done.",
                 font_size=18, color=COLORS["text"])


def create_what_mcp_is_slide(prs):
    """Slide 04: What MCP Is"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["light"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(3), Inches(0.3),
                 "TIER 1 / EXPLORER", font_size=12, bold=True, color=COLORS["gray"])

    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(5), Inches(0.4),
                 "WHAT MCP ACTUALLY IS", font_size=12, bold=True, color=COLORS["red"])

    add_text_box(slide, Inches(0.6), Inches(2), Inches(8), Inches(0.8),
                 "A protocol.", font_size=48, bold=True, color=COLORS["dark"])

    # USB-C analogy
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(10), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Like "
    p.font.size = Pt(48)
    p.font.bold = True
    run = p.add_run()
    run.text = "USB-C"
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = COLORS["red"]
    run2 = p.add_run()
    run2.text = " for AI."
    run2.font.size = Pt(48)
    run2.font.bold = True

    add_text_box(slide, Inches(0.6), Inches(5), Inches(9), Inches(1),
                 "One standard that lets any AI assistant\nconnect to any app.",
                 font_size=22, color=COLORS["text"])


def create_three_parts_slide(prs):
    """Slide 05: Three Parts"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["light"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(3), Inches(0.3),
                 "TIER 1 / EXPLORER", font_size=12, bold=True, color=COLORS["gray"])

    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(3), Inches(0.4),
                 "ARCHITECTURE", font_size=12, bold=True, color=COLORS["red"])

    add_text_box(slide, Inches(0.6), Inches(1.7), Inches(8), Inches(0.8),
                 "How it works", font_size=48, bold=True, color=COLORS["dark"])

    # Three boxes
    box_y = Inches(3)
    box_width = Inches(3)
    box_height = Inches(1.5)

    # Client box
    client_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), box_y, box_width, box_height)
    set_shape_fill(client_box, RGBColor(0xFF, 0xFF, 0xFF))
    set_shape_border(client_box, COLORS["dark"], Pt(3))
    add_text_box(slide, Inches(0.6), box_y + Inches(0.3), box_width, Inches(0.3),
                 "01", font_size=10, bold=True, color=COLORS["gray"], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.6), box_y + Inches(0.6), box_width, Inches(0.5),
                 "Client", font_size=28, bold=True, color=COLORS["dark"], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.6), box_y + Inches(1.7), box_width, Inches(1.2),
                 "Claude Desktop\nChatGPT\nCursor", font_size=18, color=COLORS["text"])

    # Arrow
    add_text_box(slide, Inches(3.8), box_y + Inches(0.5), Inches(0.8), Inches(0.6),
                 "→", font_size=36, color=COLORS["gray"], align=PP_ALIGN.CENTER)

    # Server box (highlighted)
    server_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), box_y, box_width, box_height)
    set_shape_fill(server_box, COLORS["red"])
    server_box.line.fill.background()
    add_text_box(slide, Inches(4.8), box_y + Inches(0.3), box_width, Inches(0.3),
                 "02", font_size=10, bold=True, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.8), box_y + Inches(0.6), box_width, Inches(0.5),
                 "Server", font_size=28, bold=True, color=COLORS["light"], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.8), box_y + Inches(1.7), box_width, Inches(1.2),
                 "Miro MCP Server\nPlaywright MCP Server\nGLEIF MCP Server", font_size=18, color=COLORS["text"])

    # Arrow
    add_text_box(slide, Inches(8), box_y + Inches(0.5), Inches(0.8), Inches(0.6),
                 "→", font_size=36, color=COLORS["gray"], align=PP_ALIGN.CENTER)

    # App box
    app_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9), box_y, box_width, box_height)
    set_shape_fill(app_box, RGBColor(0xFF, 0xFF, 0xFF))
    set_shape_border(app_box, COLORS["dark"], Pt(3))
    add_text_box(slide, Inches(9), box_y + Inches(0.3), box_width, Inches(0.3),
                 "03", font_size=10, bold=True, color=COLORS["gray"], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(9), box_y + Inches(0.6), box_width, Inches(0.5),
                 "App", font_size=28, bold=True, color=COLORS["dark"], align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(9), box_y + Inches(1.7), box_width, Inches(1.2),
                 "Miro\nYour browser\nGLEIF API", font_size=18, color=COLORS["text"])


def create_what_servers_provide_slide(prs):
    """Slide 06: What Servers Provide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["light"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(3), Inches(0.3),
                 "TIER 1 / EXPLORER", font_size=12, bold=True, color=COLORS["gray"])

    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(3), Inches(0.4),
                 "ARCHITECTURE", font_size=12, bold=True, color=COLORS["red"])

    add_text_box(slide, Inches(0.6), Inches(1.7), Inches(10), Inches(0.8),
                 "What servers provide", font_size=48, bold=True, color=COLORS["dark"])

    # Three columns for Tools, Resources, Prompts
    col_width = Inches(3.8)
    col_y = Inches(2.8)

    for i, (title, subtitle, desc, examples) in enumerate([
        ("Tools", "Actions", "Things the AI can do", ["Search for a company", "Create a board", "Take a screenshot"]),
        ("Resources", "Data", "Things the AI can read", ["Files on your computer", "Database records", "Live API data"]),
        ("Prompts", "Templates", "Reusable instructions", ["Code review checklist", "Summary format", "Report template"])
    ]):
        x = Inches(0.6) + i * col_width

        # Box
        is_first = i == 0
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, col_y, Inches(3.2), Inches(1.2))
        if is_first:
            set_shape_fill(box, COLORS["red"])
            box.line.fill.background()
            num_color = COLORS["muted"]
            title_color = COLORS["light"]
        else:
            set_shape_fill(box, RGBColor(0xFF, 0xFF, 0xFF))
            set_shape_border(box, COLORS["dark"], Pt(3))
            num_color = COLORS["gray"]
            title_color = COLORS["dark"]

        add_text_box(slide, x, col_y + Inches(0.2), Inches(3.2), Inches(0.3),
                     f"0{i+1}", font_size=10, bold=True, color=num_color, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, col_y + Inches(0.5), Inches(3.2), Inches(0.5),
                     title, font_size=28, bold=True, color=title_color, align=PP_ALIGN.CENTER)

        # Subtitle and description
        add_text_box(slide, x, col_y + Inches(1.4), Inches(3.2), Inches(0.4),
                     subtitle, font_size=18, bold=True, color=COLORS["dark"])
        add_text_box(slide, x, col_y + Inches(1.8), Inches(3.2), Inches(0.4),
                     desc, font_size=16, color=COLORS["text"])

        # Examples
        examples_text = "\n".join([f"→ {e}" for e in examples])
        add_text_box(slide, x, col_y + Inches(2.3), Inches(3.2), Inches(1.5),
                     examples_text, font_size=18, color=COLORS["text"])


def create_cloud_vs_local_slide(prs):
    """Slide 07: Cloud vs Local"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["light"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(3), Inches(0.3),
                 "TIER 1 / EXPLORER", font_size=12, bold=True, color=COLORS["gray"])

    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(3), Inches(0.4),
                 "ARCHITECTURE", font_size=12, bold=True, color=COLORS["red"])

    add_text_box(slide, Inches(0.6), Inches(1.7), Inches(10), Inches(0.8),
                 "Apps can be anywhere.", font_size=48, bold=True, color=COLORS["dark"])

    # Cloud Apps column
    add_text_box(slide, Inches(0.6), Inches(2.8), Inches(5), Inches(0.5),
                 "Cloud Apps", font_size=24, bold=True, color=COLORS["gray"])
    add_text_box(slide, Inches(0.6), Inches(3.4), Inches(5), Inches(2.5),
                 "Miro\nProductPlan\nGLEIF API\nJira, Slack, Notion...",
                 font_size=20, color=COLORS["text_dark"])

    # Local Apps column
    add_text_box(slide, Inches(7), Inches(2.8), Inches(5), Inches(0.5),
                 "Local Apps", font_size=24, bold=True, color=COLORS["red"])
    add_text_box(slide, Inches(7), Inches(3.4), Inches(5), Inches(2.5),
                 "Your browser\nFiles on your machine\nLocal git repos\nInstalled apps",
                 font_size=20, color=COLORS["text_dark"])

    # Insight
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(6), Inches(10), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "MCP works with both. "
    p.font.size = Pt(20)
    p.font.bold = True
    run = p.add_run()
    run.text = "Data stays where you want it."
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = COLORS["red"]


def create_tier2_intro_slide(prs):
    """Slide 08: Tier 2 Intro"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["dark"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(2.5), Inches(3), Inches(0.3),
                 "TIER 2 / USER", font_size=12, bold=True, color=COLORS["gray"])

    # Title with highlight
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3), Inches(10), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "How to "
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS["light"]
    run = p.add_run()
    run.text = "Use"
    run.font.size = Pt(72)
    run.font.bold = True
    run.font.color.rgb = COLORS["red"]
    run2 = p.add_run()
    run2.text = " MCPs"
    run2.font.size = Pt(72)
    run2.font.bold = True
    run2.font.color.rgb = COLORS["light"]

    add_text_box(slide, Inches(0.6), Inches(4.8), Inches(8), Inches(1),
                 "Can follow instructions, no coding required.\nThree live demos.",
                 font_size=24, color=COLORS["muted"])


def create_get_claude_desktop_slide(prs):
    """Slide 08b: Get Claude Desktop"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["light"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(3), Inches(0.3),
                 "TIER 2 / USER", font_size=12, bold=True, color=COLORS["gray"])

    add_text_box(slide, Inches(0.6), Inches(1.5), Inches(10), Inches(0.8),
                 "Get Claude Desktop", font_size=48, bold=True, color=COLORS["dark"])
    add_text_box(slide, Inches(0.6), Inches(2.2), Inches(5), Inches(0.4),
                 "claude.ai/download", font_size=20, bold=True, color=COLORS["red"])

    # Mac column
    add_text_box(slide, Inches(0.6), Inches(3.2), Inches(5), Inches(0.4),
                 "MAC", font_size=14, bold=True, color=COLORS["gray"])
    add_text_box(slide, Inches(0.6), Inches(3.7), Inches(5), Inches(2),
                 "→ Download the .dmg file\n→ Drag Claude to Applications\n→ Open and sign in",
                 font_size=20, color=COLORS["text_dark"])

    # Windows column
    add_text_box(slide, Inches(7), Inches(3.2), Inches(5), Inches(0.4),
                 "WINDOWS", font_size=14, bold=True, color=COLORS["gray"])
    add_text_box(slide, Inches(7), Inches(3.7), Inches(5), Inches(2),
                 "→ Download and run installer\n→ Follow the wizard\n→ Open and sign in",
                 font_size=20, color=COLORS["text_dark"])


def create_demo_slide(prs, demo_num, title, demo_type, prompts, points):
    """Generic demo slide template"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["light"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(3), Inches(0.3),
                 "TIER 2 / USER", font_size=12, bold=True, color=COLORS["gray"])

    # Demo badge
    badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.3), Inches(1.2), Inches(0.4))
    set_shape_fill(badge, COLORS["red"])
    badge.line.fill.background()
    add_text_box(slide, Inches(0.6), Inches(1.35), Inches(1.2), Inches(0.35),
                 f"DEMO {demo_num}", font_size=12, bold=True, color=COLORS["light"], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.6), Inches(1.9), Inches(8), Inches(0.8),
                 title, font_size=48, bold=True, color=COLORS["dark"])
    add_text_box(slide, Inches(0.6), Inches(2.6), Inches(5), Inches(0.4),
                 demo_type, font_size=18, bold=True, color=COLORS["gray"])

    # Prompts column
    add_text_box(slide, Inches(0.6), Inches(3.3), Inches(5), Inches(0.4),
                 "LIVE PROMPTS", font_size=14, bold=True, color=COLORS["gray"])

    prompt_y = Inches(3.8)
    for prompt in prompts:
        prompt_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), prompt_y, Inches(5.5), Inches(0.7))
        set_shape_fill(prompt_box, COLORS["code_bg"])
        prompt_box.line.fill.background()
        # Left accent
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), prompt_y, Inches(0.08), Inches(0.7))
        set_shape_fill(accent, COLORS["red"])
        accent.line.fill.background()
        add_text_box(slide, Inches(0.8), prompt_y + Inches(0.15), Inches(5.2), Inches(0.5),
                     prompt, font_size=14, color=COLORS["dark"])
        prompt_y += Inches(0.85)

    # Points column
    add_text_box(slide, Inches(7), Inches(3.3), Inches(5), Inches(0.4),
                 "UNDER THE HOOD" if demo_num == 1 else "THE DIFFERENCE" if demo_num == 2 else "WHAT IT TAKES",
                 font_size=14, bold=True, color=COLORS["gray"])

    points_text = "\n".join([f"→ {p}" for p in points])
    add_text_box(slide, Inches(7), Inches(3.8), Inches(5.5), Inches(3),
                 points_text, font_size=18, color=COLORS["text_dark"])


def create_install_gleif_slide(prs):
    """Slide 11b: Install GLEIF"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["light"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(3), Inches(0.3),
                 "TIER 2 / USER", font_size=12, bold=True, color=COLORS["gray"])

    add_text_box(slide, Inches(0.6), Inches(1.5), Inches(10), Inches(0.8),
                 "Try it: Add GLEIF", font_size=48, bold=True, color=COLORS["dark"])
    add_text_box(slide, Inches(0.6), Inches(2.2), Inches(10), Inches(0.4),
                 "github.com/olgasafonova/gleif-mcp-server/releases", font_size=18, color=COLORS["gray"])

    # Steps
    add_text_box(slide, Inches(0.6), Inches(3.2), Inches(5), Inches(0.4),
                 "STEPS", font_size=14, bold=True, color=COLORS["gray"])
    add_text_box(slide, Inches(0.6), Inches(3.7), Inches(5.5), Inches(2.5),
                 "→ Download the binary for your system\n→ Open config folder\n→ Add the config (see GitHub README)\n→ Restart Claude Desktop",
                 font_size=20, color=COLORS["text_dark"])

    # Config folder
    add_text_box(slide, Inches(7), Inches(3.2), Inches(5), Inches(0.4),
                 "CONFIG FOLDER", font_size=14, bold=True, color=COLORS["gray"])
    add_text_box(slide, Inches(7), Inches(3.7), Inches(5.5), Inches(1.5),
                 "→ Mac: Cmd+Shift+G, paste path\n→ Windows: Win+R, paste path",
                 font_size=20, color=COLORS["text_dark"])


def create_tier3_intro_slide(prs):
    """Slide 12: Tier 3 Intro"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["dark"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(2.5), Inches(3), Inches(0.3),
                 "TIER 3 / ARCHITECT", font_size=12, bold=True, color=COLORS["gray"])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3), Inches(10), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "How to "
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS["light"]
    run = p.add_run()
    run.text = "Build"
    run.font.size = Pt(72)
    run.font.bold = True
    run.font.color.rgb = COLORS["red"]
    run2 = p.add_run()
    run2.text = " MCPs"
    run2.font.size = Pt(72)
    run2.font.bold = True
    run2.font.color.rgb = COLORS["light"]

    add_text_box(slide, Inches(0.6), Inches(4.8), Inches(8), Inches(1),
                 "For developers who want to create their own.\nReal patterns from production servers.",
                 font_size=24, color=COLORS["muted"])


def create_architecture_slide(prs):
    """Slide 13: Architecture Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["light"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(3), Inches(0.3),
                 "TIER 3 / ARCHITECT", font_size=12, bold=True, color=COLORS["gray"])

    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(3), Inches(0.4),
                 "STRUCTURE", font_size=12, bold=True, color=COLORS["red"])

    add_text_box(slide, Inches(0.6), Inches(1.7), Inches(10), Inches(0.8),
                 "Architecture Overview", font_size=48, bold=True, color=COLORS["dark"])

    # Code block (dark background)
    code_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.8), Inches(5.5), Inches(4))
    set_shape_fill(code_box, COLORS["dark"])
    code_box.line.fill.background()

    code_text = """server/
├── main.go         # Entry point
├── tools/
│   ├── definitions.go  # What can it do?
│   └── handlers.go     # How does it do it?
└── client/
    └── client.go       # API client"""

    add_text_box(slide, Inches(0.8), Inches(3), Inches(5), Inches(3.5),
                 code_text, font_size=16, color=COLORS["light"])

    # Explanations
    explanations = [
        ("main.go", "Entry point. Registers tools, starts the server."),
        ("definitions.go", "Tool names, descriptions, parameter schemas. What Claude sees."),
        ("handlers.go", "The actual implementation. Calls APIs, processes data, returns results."),
        ("client.go", "HTTP client with auth, rate limiting, caching. Talks to the real API."),
    ]

    y = Inches(2.8)
    for title, desc in explanations:
        add_text_box(slide, Inches(7), y, Inches(5.5), Inches(0.4),
                     title, font_size=18, bold=True, color=COLORS["red"])
        add_text_box(slide, Inches(7), y + Inches(0.35), Inches(5.5), Inches(0.6),
                     desc, font_size=16, color=COLORS["text"])
        y += Inches(0.95)


def create_patterns_slide(prs):
    """Slide 14: Production Patterns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["light"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(3), Inches(0.3),
                 "TIER 3 / ARCHITECT", font_size=12, bold=True, color=COLORS["gray"])

    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(3), Inches(0.4),
                 "LESSONS LEARNED", font_size=12, bold=True, color=COLORS["red"])

    add_text_box(slide, Inches(0.6), Inches(1.7), Inches(10), Inches(0.8),
                 "Production Patterns", font_size=48, bold=True, color=COLORS["dark"])

    patterns = [
        ("Error Handling", "Wrap API errors with context. Catch failures gracefully. Return clear error messages."),
        ("Authentication", "Store API keys safely, not in code. Check they exist at start."),
        ("Caching", "Remember recent results. Fewer API calls, faster responses."),
        ("Rate Limiting", "Don't hit APIs too fast. Slow down when needed."),
        ("Bulk Operations", "Group similar requests. One batch beats many singles."),
        ("Retries", "Retry temporary failures. Skip bad requests. Wait longer between attempts."),
    ]

    col_width = Inches(5.8)
    for i, (title, desc) in enumerate(patterns):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + col * col_width
        y = Inches(2.8) + row * Inches(1.3)

        add_text_box(slide, x, y, Inches(5), Inches(0.4),
                     title, font_size=18, bold=True, color=COLORS["dark"])
        add_text_box(slide, x, y + Inches(0.35), Inches(5.5), Inches(0.8),
                     desc, font_size=16, color=COLORS["text"])


def create_getting_started_slide(prs):
    """Slide 15: Getting Started"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["light"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(3), Inches(0.3),
                 "TIER 3 / ARCHITECT", font_size=12, bold=True, color=COLORS["gray"])

    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(5), Inches(0.4),
                 "YOUR FIRST MCP SERVER", font_size=12, bold=True, color=COLORS["red"])

    add_text_box(slide, Inches(0.6), Inches(1.7), Inches(10), Inches(0.8),
                 "Getting Started", font_size=48, bold=True, color=COLORS["dark"])

    options = [
        ("Fork GLEIF MCP", "Simplest, 12 tools", "Clean, minimal codebase. Read-only API. Perfect for learning.", True),
        ("Fork Miro MCP", "Most advanced, 77 tools", "Full CRUD operations, OAuth flow, bulk actions.", False),
        ("Start fresh", "Go, Python, or TypeScript SDK", "Official SDKs handle protocol. You focus on tools.", False),
    ]

    for i, (title, subtitle, desc, recommended) in enumerate(options):
        x = Inches(0.6) + i * Inches(4)

        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.8), Inches(3.6), Inches(2.2))
        set_shape_fill(box, RGBColor(0xFF, 0xFF, 0xFF))
        if recommended:
            set_shape_border(box, COLORS["red"], Pt(3))
            # Recommended badge
            badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.2), Inches(2.95), Inches(1.5), Inches(0.35))
            set_shape_fill(badge, COLORS["red"])
            badge.line.fill.background()
            add_text_box(slide, x + Inches(0.2), Inches(2.98), Inches(1.5), Inches(0.3),
                         "RECOMMENDED", font_size=10, bold=True, color=COLORS["light"], align=PP_ALIGN.CENTER)
        else:
            set_shape_border(box, COLORS["border"], Pt(2))

        add_text_box(slide, x + Inches(0.2), Inches(3.4) if recommended else Inches(3.1), Inches(3.2), Inches(0.5),
                     title, font_size=22, bold=True, color=COLORS["dark"])
        add_text_box(slide, x + Inches(0.2), Inches(3.85) if recommended else Inches(3.55), Inches(3.2), Inches(0.35),
                     subtitle, font_size=14, color=COLORS["gray"])
        add_text_box(slide, x + Inches(0.2), Inches(4.25) if recommended else Inches(3.95), Inches(3.2), Inches(0.8),
                     desc, font_size=14, color=COLORS["text"])

    # SDK note
    add_text_box(slide, Inches(0.6), Inches(5.3), Inches(12), Inches(1.5),
                 "Python and TypeScript are the most popular. Why Go? Single binary, no runtime dependencies, starts instantly.\nYou decide what works for you.",
                 font_size=18, color=COLORS["text"])


def create_distilled_slide(prs):
    """Slide 16: Distilled/Takeaways"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["light"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(3), Inches(0.4),
                 "WRAP-UP", font_size=12, bold=True, color=COLORS["red"])

    add_text_box(slide, Inches(0.6), Inches(1.7), Inches(10), Inches(0.8),
                 "Distilled", font_size=48, bold=True, color=COLORS["dark"])

    takeaways = [
        ("Explorers", "MCP lets AI ", "take action", ", not just answer questions"),
        ("Users", "Setup is simple: ", "download, configure, restart", ""),
        ("Architects", "Start small. Build for ", "real value", ""),
    ]

    y = Inches(3)
    for tier, prefix, highlight, suffix in takeaways:
        add_text_box(slide, Inches(0.6), y, Inches(1.5), Inches(0.5),
                     tier.upper(), font_size=14, bold=True, color=COLORS["gray"])

        txBox = slide.shapes.add_textbox(Inches(2.2), y, Inches(10), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = prefix
        p.font.size = Pt(26)
        p.font.bold = True
        run = p.add_run()
        run.text = highlight
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = COLORS["red"]
        if suffix:
            run2 = p.add_run()
            run2.text = suffix
            run2.font.size = Pt(26)
            run2.font.bold = True

        y += Inches(1)


def create_testimonial_slide(prs):
    """Slide 17: Testimonials"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["light"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(3), Inches(0.4),
                 "REAL VALUE", font_size=12, bold=True, color=COLORS["red"])

    # Miro section
    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(2), Inches(0.5),
                 "Miro", font_size=28, bold=True, color=COLORS["dark"])

    add_text_box(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(1.5),
                 '"Impressive work on Miro MCP, thank you!"\n— Andrey Khusid, CEO & Co-founder, Miro\n\n"You are ahead of the curve."\n— David Ross, Chief Product Evangelist, Miro',
                 font_size=16, color=COLORS["text"])

    # Divider
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1), Inches(0.02), Inches(5.5))
    set_shape_fill(divider, COLORS["border"])
    divider.line.fill.background()

    # Colleague quotes
    add_text_box(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(0.5),
                 "Tietoevry", font_size=28, bold=True, color=COLORS["dark"])

    quotes = [
        ('"SO useful in making everyday work more efficient"', "Jenny Felldin, Head of Product"),
        ('"A golden nugget to help sales"', "Jenny Felldin, Head of Product"),
        ('"This might be one of the best gifts ever"', "Ari Räisänen, Lead Software Architect"),
        ('"Design Team is super excited"', "Aleksandra Bratek, UI Designer"),
        ('"Love it!"', "Morten Lerfald, Senior Business Developer"),
    ]

    y = Inches(1.8)
    for quote, author in quotes:
        add_text_box(slide, Inches(7), y, Inches(5.5), Inches(0.5),
                     quote, font_size=14, bold=True, color=COLORS["dark"])
        add_text_box(slide, Inches(7), y + Inches(0.35), Inches(5.5), Inches(0.35),
                     f"— {author}", font_size=12, color=COLORS["gray"])
        y += Inches(0.85)


def create_resources_slide(prs):
    """Slide 18: Resources"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["dark"])
    background.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(3), Inches(0.4),
                 "RESOURCES", font_size=12, bold=True, color=COLORS["red"])

    add_text_box(slide, Inches(0.6), Inches(1.8), Inches(10), Inches(0.8),
                 "Learn More", font_size=36, bold=True, color=COLORS["light"])

    resources = [
        ("MCP Specification", "modelcontextprotocol.io"),
        ("Go SDK", "github.com/modelcontextprotocol/go-sdk"),
        ("TypeScript SDK", "github.com/modelcontextprotocol/typescript-sdk"),
        ("Python SDK", "github.com/modelcontextprotocol/python-sdk"),
        ("MCP servers I built", "github.com/olgasafonova"),
    ]

    y = Inches(2.8)
    for title, url in resources:
        txBox = slide.shapes.add_textbox(Inches(0.6), y, Inches(10), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS["light"]
        run = p.add_run()
        run.text = f"  {url}"
        run.font.size = Pt(20)
        run.font.bold = False
        run.font.color.rgb = COLORS["muted"]
        y += Inches(0.6)


def create_thankyou_slide(prs):
    """Slide 19: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(background, COLORS["dark"])
    background.line.fill.background()

    add_text_box(slide, Inches(4), Inches(1.5), Inches(5.333), Inches(1),
                 "Thank you!", font_size=64, bold=True, color=COLORS["red"], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(4), Inches(2.8), Inches(5.333), Inches(0.6),
                 "Olga Safonova", font_size=36, bold=True, color=COLORS["light"], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(3), Inches(3.5), Inches(7.333), Inches(1),
                 "AI Product Leader | Building AI tools in Go\nEx-Workday, Siteimprove, Saxo Bank | Keynote Speaker",
                 font_size=18, color=COLORS["gray"], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(4), Inches(5), Inches(5.333), Inches(0.6),
                 "Connect", font_size=32, bold=True, color=COLORS["text"], align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(5.8), Inches(9.333), Inches(0.5),
                 "LinkedIn                    GitHub                    Substack",
                 font_size=16, bold=True, color=COLORS["gray"], align=PP_ALIGN.CENTER)


def main():
    """Generate the complete presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Generating slides...")

    # All slides in order
    create_title_slide(prs)                    # 01
    create_problem_slide(prs)                  # 02
    create_rag_vs_mcp_slide(prs)              # 03
    create_what_mcp_is_slide(prs)             # 04
    create_three_parts_slide(prs)             # 05
    create_what_servers_provide_slide(prs)    # 06
    create_cloud_vs_local_slide(prs)          # 07
    create_tier2_intro_slide(prs)             # 08
    create_get_claude_desktop_slide(prs)      # 08b

    # Demo slides
    create_demo_slide(prs, 1, "GLEIF", "Cloud App Pattern",
                     ['"Use the GLEIF server to look up Apple Inc\'s LEI"',
                      '"Who is Apple\'s ultimate parent company?"'],
                     ["Data lives in the cloud (GLEIF API)", "Lookup, validate, trace ownership",
                      "Structured output, not web scraping", "Real financial identifiers (LEI)"])

    create_demo_slide(prs, 2, "Playwright", "Local App Pattern",
                     ['"Open a browser and go to producthunt.com"',
                      '"Take a screenshot of the page"'],
                     ["Your browser, your machine", "Navigate, fill forms, capture pages",
                      "Privacy: data stays on your machine", "No API keys needed for local apps"])

    create_demo_slide(prs, 3, "Configuration", "What you actually do",
                     ['// claude_desktop_config.json'],
                     ["Edit one file", "Give it a name, tell it how to start",
                      "Restart to activate", "No coding to use an MCP"])

    create_install_gleif_slide(prs)           # 11b
    create_tier3_intro_slide(prs)             # 12
    create_architecture_slide(prs)            # 13
    create_patterns_slide(prs)                # 14
    create_getting_started_slide(prs)         # 15
    create_distilled_slide(prs)               # 16
    create_testimonial_slide(prs)             # 17
    create_resources_slide(prs)               # 18
    create_thankyou_slide(prs)                # 19

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "mcp-workshop.pptx")
    prs.save(output_path)
    print(f"✓ Saved: {output_path}")
    print(f"  {len(prs.slides)} slides generated")


if __name__ == "__main__":
    main()
