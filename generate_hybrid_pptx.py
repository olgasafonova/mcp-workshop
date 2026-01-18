#!/usr/bin/env python3
"""
Generate hybrid PPTX: screenshots as backgrounds + editable text overlays.
Reads metadata from slide-screenshots/metadata.json
"""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Brand colors
COLORS = {
    "red": RGBColor(0xE5, 0x39, 0x35),
    "dark": RGBColor(0x1A, 0x1A, 0x1A),
    "light": RGBColor(0xFA, 0xFA, 0xFA),
}

# Dark slides (text should be light colored)
DARK_SLIDES = [
    'slide-08-tier2-intro.html',
    'slide-12-tier3-intro.html',
    'slide-18-resources.html',
    'slide-19-thankyou.html',
]


def add_text_overlay(slide, overlay, is_dark_slide=False):
    """Add an editable text box overlay."""
    # Convert percentages to inches
    left = SLIDE_WIDTH * (overlay['left'] / 100)
    top = SLIDE_HEIGHT * (overlay['top'] / 100)
    width = SLIDE_WIDTH * (overlay['width'] / 100)
    height = SLIDE_HEIGHT * (overlay['height'] / 100)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = overlay['text']
    p.font.size = Pt(overlay.get('fontSize', 24))
    p.font.bold = overlay.get('bold', False)
    p.font.name = 'Arial'

    # Handle alignment
    align = overlay.get('align', 'left')
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT

    # Set text color based on slide background
    # Check if text contains highlight words
    text = overlay['text']
    if is_dark_slide:
        # Light text on dark background
        if 'Thank you' in text:
            p.font.color.rgb = COLORS["red"]
        else:
            p.font.color.rgb = COLORS["light"]
    else:
        p.font.color.rgb = COLORS["dark"]

    # Make text box background transparent
    txBox.fill.background()

    return txBox


def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    screenshots_dir = os.path.join(script_dir, 'slide-screenshots')
    metadata_path = os.path.join(screenshots_dir, 'metadata.json')

    if not os.path.exists(metadata_path):
        print(f"Error: {metadata_path} not found.")
        print("Run 'node generate_hybrid_pptx.js' first to create screenshots.")
        return

    with open(metadata_path) as f:
        slides_data = json.load(f)

    print(f"Creating PPTX from {len(slides_data)} slides...")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add screenshot as background image (full bleed)
        img_path = slide_data['path']
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        else:
            print(f"  ⚠ Missing: {img_path}")
            continue

        # Determine if this is a dark slide
        is_dark = slide_data['file'] in DARK_SLIDES

        # Add editable text overlays
        for overlay in slide_data.get('overlays', []):
            add_text_overlay(slide, overlay, is_dark)

        print(f"  ✓ Slide {slide_data['index'] + 1}: {slide_data['file']}")

    output_path = os.path.join(script_dir, 'mcp-workshop-hybrid.pptx')
    prs.save(output_path)
    print(f"\n✓ Saved: {output_path}")
    print(f"  {len(prs.slides)} slides with editable text overlays")


if __name__ == '__main__':
    main()
